"""Generate DATARADAR Strategic Overview deck (14 slides, current state)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Palette
BG = RGBColor(0x0B, 0x14, 0x20)
FG = RGBColor(0xF0, 0xF3, 0xF7)
BLUE = RGBColor(0x4A, 0x9E, 0xFF)
GREEN = RGBColor(0x30, 0xD1, 0x58)
GOLD = RGBColor(0xC8, 0xA3, 0x55)
AMBER = RGBColor(0xFF, 0xB3, 0x30)
RED = RGBColor(0xFF, 0x5A, 0x5A)
MUTED = RGBColor(0x50, 0x5A, 0x6A)
DIM = RGBColor(0x7A, 0x84, 0x94)
PANEL = RGBColor(0x12, 0x1E, 0x2F)
PANEL_DARK = RGBColor(0x08, 0x10, 0x1C)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


def set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=FG, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rich_lines(slide, left, top, width, height, lines, size=18,
                   color=FG, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                   line_spacing=1.25, font=FONT):
    """lines: list of (text, bold, color_or_None, size_or_None)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, spec in enumerate(lines):
        text, bold, c, sz = spec
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(sz or size)
        run.font.bold = bold
        run.font.color.rgb = c or color
    return tb


def wordmark(slide):
    add_text(slide, Inches(11.7), Inches(7.15), Inches(1.5), Inches(0.3),
             "DATARADAR", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def accent_bar(slide, color=BLUE, top=None, left=None, w=None, h=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left or Inches(0.6), top or Inches(1.2),
        w or Inches(0.12), h or Inches(0.7))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def title(slide, text, accent=BLUE, y=Inches(0.55)):
    accent_bar(slide, color=accent, top=y, left=Inches(0.6),
               w=Inches(0.14), h=Inches(0.8))
    add_text(slide, Inches(0.9), y, Inches(12), Inches(0.9),
             text, size=36, bold=True, color=FG)


def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    return s


def rounded_box(slide, left, top, w, h, fill=PANEL, line=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(1)
    box.shadow.inherit = False
    box.text_frame.text = ""
    return box


# =========================================================================
# Slide 1 — Title
# =========================================================================
s = new_slide()
big = s.shapes.add_shape(MSO_SHAPE.OVAL,
                         Inches(9.2), Inches(-2.5),
                         Inches(7.5), Inches(7.5))
big.line.fill.background()
big.fill.solid()
big.fill.fore_color.rgb = RGBColor(0x10, 0x2A, 0x4D)
big.shadow.inherit = False

bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                         Inches(0), Inches(3.5),
                         Inches(5), Inches(0.08))
bar.line.fill.background()
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE

dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                         Inches(4.85), Inches(3.42),
                         Inches(0.24), Inches(0.24))
dot.line.fill.background()
dot.fill.solid()
dot.fill.fore_color.rgb = GOLD

# Small triangle accent
tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                         Inches(0.8), Inches(5.55),
                         Inches(0.35), Inches(0.35))
tri.line.fill.background()
tri.fill.solid()
tri.fill.fore_color.rgb = GREEN

add_text(s, Inches(0.8), Inches(2.2), Inches(10), Inches(1.4),
         "DATARADAR", size=84, bold=True, color=FG)
add_text(s, Inches(0.8), Inches(3.7), Inches(10), Inches(0.7),
         "Data-driven eBay inventory intelligence",
         size=24, color=BLUE)
add_text(s, Inches(0.8), Inches(4.55), Inches(10), Inches(0.5),
         "JJ Shay  ·  Gauntlet Gallery  ·  2026",
         size=18, color=FG)
add_text(s, Inches(0.8), Inches(6.4), Inches(10), Inches(0.4),
         "web-production-15df7.up.railway.app",
         size=14, color=GOLD)
wordmark(s)

# =========================================================================
# Slide 2 — The Problem
# =========================================================================
s = new_slide()
title(s, "eBay resellers are flying blind", accent=GOLD)

bullets = [
    ("•  54,000+ comparable sales exist on the web — fragmented, unindexed", False, FG, 20),
    ("•  Sellers price by gut or by the single lowest competitor (race to bottom)", False, FG, 20),
    ("•  No systematic way to ask: \"what should this actually sell for?\"", False, FG, 20),
    ("•  Key-date events move prices 5–35% — most sellers miss the window", False, FG, 20),
]
add_rich_lines(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(3.8),
               bullets, line_spacing=1.5)

rounded_box(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.9),
            fill=RGBColor(0x15, 0x24, 0x3A))
add_text(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.9),
         "Price accuracy = 60%+ of net margin",
         size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
wordmark(s)

# =========================================================================
# Slide 3 — The Solution
# =========================================================================
s = new_slide()
title(s, "DATARADAR in one sentence", accent=BLUE)

add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.6),
         "Every eBay listing, priced by 54,000 comps\nand 4 AI models, in real time.",
         size=32, bold=True, color=FG, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# Flow diagram: 6 boxes
steps = ["eBay\nInventory", "54k Comp\nDB", "Smart\nAnchor",
         "4-LLM\nReview", "Consensus\nRange", "One-Click\nReprice"]
colors = [BLUE, GOLD, BLUE, GREEN, GOLD, BLUE]
n = len(steps)
total_w = Inches(12.3)
box_w = Inches(1.75)
gap = (total_w - box_w * n) / (n - 1)
start_x = Inches(0.5)
row_y = Inches(5.1)
box_h = Inches(1.3)

centers = []
for i, (label, col) in enumerate(zip(steps, colors)):
    x = start_x + (box_w + gap) * i
    box = rounded_box(s, x, row_y, box_w, box_h,
                      fill=PANEL, line=col)
    tf = box.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = FONT
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = FG
    centers.append((x + box_w, row_y + box_h / 2))

for i in range(n - 1):
    x_end, y_mid = centers[i]
    next_x = start_x + (box_w + gap) * (i + 1)
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               x_end + Emu(20000),
                               y_mid - Inches(0.12),
                               next_x - x_end - Emu(40000),
                               Inches(0.24))
    arrow.line.fill.background()
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BLUE
wordmark(s)

# =========================================================================
# Slide 4 — Data Moat
# =========================================================================
s = new_slide()
title(s, "The 54,000-record comp database", accent=GOLD)

bullets = [
    ("•  Deduplicated across Shepard Fairey, Death NYC, KAWS, Banksy, Mr. Brainwash, Bearbrick", False, FG, 17),
    ("•  Indexed: artist, work_id, colorway, medium, edition size, signed/numbered status", False, FG, 17),
    ("•  Weekly scraper pipeline + manual curation loop (operator swipes L=reject R=approve)", False, FG, 17),
    ("•  Powers /prices mobile UI for on-the-go lookups", False, FG, 17),
]
add_rich_lines(s, Inches(0.9), Inches(1.85), Inches(8.0), Inches(3.6),
               bullets, line_spacing=1.5)

# Three stat chips across the bottom
stats = [
    ("54,000+", "comps in DB", GOLD),
    ("<200ms", "median search", BLUE),
    ("99%", "dedup rate", GREEN),
]
chip_w = Inches(3.7)
chip_h = Inches(1.5)
gap_x = Inches(0.25)
total = chip_w * 3 + gap_x * 2
start = (SW - total) / 2
cy = Inches(5.5)
for i, (big_n, lbl, col) in enumerate(stats):
    x = start + (chip_w + gap_x) * i
    rounded_box(s, x, cy, chip_w, chip_h, fill=PANEL, line=col)
    add_text(s, x, cy + Inches(0.15), chip_w, Inches(0.9),
             big_n, size=40, bold=True, color=col,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, cy + Inches(1.05), chip_w, Inches(0.4),
             lbl, size=13, color=DIM, align=PP_ALIGN.CENTER)
wordmark(s)

# =========================================================================
# Slide 5 — 4-LLM Consensus Review with Price Ranges
# =========================================================================
s = new_slide()
title(s, "Four models. One range.", accent=GREEN)

models = [
    ("CLAUDE", "Anthropic  ·  claude-sonnet-4-6", BLUE),
    ("GPT-4o", "OpenAI", GREEN),
    ("GEMINI 2.5 FLASH", "Google", GOLD),
    ("GROK 3", "xAI", BLUE),
]
grid_x0 = Inches(0.9)
grid_y0 = Inches(1.65)
cell_w = Inches(5.8)
cell_h = Inches(1.1)
gap_x = Inches(0.3)
gap_y = Inches(0.18)

for i, (name, sub, col) in enumerate(models):
    r, c = divmod(i, 2)
    x = grid_x0 + (cell_w + gap_x) * c
    y = grid_y0 + (cell_h + gap_y) * r
    rounded_box(s, x, y, cell_w, cell_h, fill=PANEL, line=col)
    add_text(s, x + Inches(0.3), y + Inches(0.17), cell_w - Inches(0.6), Inches(0.5),
             name, size=20, bold=True, color=FG)
    add_text(s, x + Inches(0.3), y + Inches(0.65), cell_w - Inches(0.6), Inches(0.4),
             sub, size=11, color=DIM)

# Highlight: range output
rounded_box(s, Inches(0.9), Inches(4.35), Inches(12.0), Inches(0.75),
            fill=RGBColor(0x14, 0x2C, 0x1A), line=GREEN)
add_text(s, Inches(0.9), Inches(4.35), Inches(12.0), Inches(0.75),
         "NEW · Each model returns {low, recommended, high} — a range, not a number",
         size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

bullets = [
    ("•  Live eBay competition fed into every prompt via Browse API", False, FG, 14),
    ("•  Per-artist prompt fragments (Fairey editions, Banksy Pest Control, KAWS colorway)", False, FG, 14),
    ("•  Parallel fan-out via ThreadPoolExecutor — 5–10s wall clock per item", False, FG, 14),
    ("•  Cached per (listing_id, comp_median bucket) — invalidates only when comp_median moves ≥$10", False, FG, 14),
]
add_rich_lines(s, Inches(0.9), Inches(5.3), Inches(12), Inches(1.5),
               bullets, line_spacing=1.3)

add_text(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.35),
         "~$0.02 per item  ·  cached until comps shift",
         size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
wordmark(s)

# =========================================================================
# Slide 6 — Confidence Scoring (NEW)
# =========================================================================
s = new_slide()
title(s, "Every pricing call rated 0–100", accent=BLUE)

# Left: component weights
rounded_box(s, Inches(0.9), Inches(1.85), Inches(6.5), Inches(4.3),
            fill=PANEL, line=BLUE)
add_text(s, Inches(1.1), Inches(2.0), Inches(6.1), Inches(0.4),
         "CONFIDENCE COMPONENTS", size=12, bold=True, color=DIM)

components = [
    ("Comp density", "40", GOLD),
    ("LLM agreement (σ)", "30", BLUE),
    ("12-mo trend stability", "15", GREEN),
    ("Live eBay competition", "15", GOLD),
]
cy = Inches(2.55)
for name, pts, col in components:
    rounded_box(s, Inches(1.1), cy, Inches(6.1), Inches(0.65),
                fill=PANEL_DARK, line=None)
    add_text(s, Inches(1.3), cy, Inches(4.2), Inches(0.65),
             name, size=16, color=FG, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.5), cy, Inches(1.6), Inches(0.65),
             pts + " pts", size=16, bold=True, color=col,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    cy += Inches(0.8)

# Right: tier band + big number example
rounded_box(s, Inches(7.7), Inches(1.85), Inches(5.2), Inches(4.3),
            fill=PANEL, line=GREEN)

# Tier chips
tiers = [("HIGH 70+", GREEN), ("MED 40–69", AMBER), ("LOW <40", RED)]
tx = Inches(7.9)
ty = Inches(2.05)
for label, col in tiers:
    rounded_box(s, tx, ty, Inches(1.5), Inches(0.45), fill=PANEL_DARK, line=col)
    add_text(s, tx, ty, Inches(1.5), Inches(0.45),
             label, size=11, bold=True, color=col,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx += Inches(1.6)

add_text(s, Inches(7.9), Inches(2.75), Inches(4.9), Inches(1.2),
         "85 / 100", size=76, bold=True, color=GREEN,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(7.9), Inches(3.95), Inches(4.9), Inches(0.4),
         "HIGH confidence", size=14, color=DIM, align=PP_ALIGN.CENTER)

reasoning = [
    ("reasoning_chain:", True, DIM, 11),
    ("•  18 comps, p75 = $385", False, FG, 12),
    ("•  4-LLM σ = $14 (tight)", False, FG, 12),
    ("•  12-mo trend flat ±3%", False, FG, 12),
]
add_rich_lines(s, Inches(7.9), Inches(4.55), Inches(4.9), Inches(1.5),
               reasoning, line_spacing=1.3)

add_text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
         "Every review exposes a reasoning_chain — click \"Why?\" to see every ingredient",
         size=14, color=GOLD, align=PP_ALIGN.CENTER)
wordmark(s)

# =========================================================================
# Slide 7 — Swipe Mode UX
# =========================================================================
s = new_slide()
title(s, "Tinder for inventory", accent=BLUE)

bullets = [
    ("•  Full-screen card view, one listing at a time", False, FG, 17),
    ("•  Three tabs per card: Comps · LLM Consensus · Train Comps (NEW)", False, FG, 17),
    ("•  Touch swipe, keyboard arrows, Esc to close", False, FG, 17),
    ("•  \"Match Consensus\" auto-pushes price to eBay in one click", False, FG, 17),
    ("•  Train Comps tab: swipe individual comps to curate the dataset", False, FG, 17),
]
add_rich_lines(s, Inches(0.9), Inches(1.85), Inches(7.7), Inches(4.8),
               bullets, line_spacing=1.45)

# Right: card sketch — 4 labeled rectangles
card_x = Inches(9.0)
card_y = Inches(1.75)
card_w = Inches(3.6)
rounded_box(s, card_x, card_y, card_w, Inches(5.2),
            fill=PANEL, line=BLUE)

# HEADER
rounded_box(s, card_x + Inches(0.2), card_y + Inches(0.2),
            card_w - Inches(0.4), Inches(0.85),
            fill=RGBColor(0x1C, 0x2B, 0x42))
add_text(s, card_x + Inches(0.2), card_y + Inches(0.2),
         card_w - Inches(0.4), Inches(0.85),
         "HEADER\nTitle · Artist · Price",
         size=11, bold=True, color=FG, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# TABS
rounded_box(s, card_x + Inches(0.2), card_y + Inches(1.2),
            card_w - Inches(0.4), Inches(0.5),
            fill=PANEL_DARK)
add_text(s, card_x + Inches(0.2), card_y + Inches(1.2),
         card_w - Inches(0.4), Inches(0.5),
         "TABS  ·  Comps | LLM | Train",
         size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# BODY
rounded_box(s, card_x + Inches(0.2), card_y + Inches(1.85),
            card_w - Inches(0.4), Inches(2.4),
            fill=RGBColor(0x1C, 0x2B, 0x42))
add_text(s, card_x + Inches(0.2), card_y + Inches(1.85),
         card_w - Inches(0.4), Inches(2.4),
         "BODY\nhistogram · year trend\n4 LLM range cards\ncomp curation stack",
         size=10, bold=True, color=FG, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# FOOTER
rounded_box(s, card_x + Inches(0.2), card_y + Inches(4.4),
            card_w - Inches(0.4), Inches(0.65),
            fill=RGBColor(0x14, 0x38, 0x24))
add_text(s, card_x + Inches(0.2), card_y + Inches(4.4),
         card_w - Inches(0.4), Inches(0.65),
         "FOOTER · Match Consensus",
         size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
wordmark(s)

# =========================================================================
# Slide 8 — Opportunities Dashboard (NEW)
# =========================================================================
s = new_slide()
title(s, "Top pricing upsides, ranked", accent=GREEN)

add_text(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(0.5),
         "GET /api/opportunities — top 10 underpriced items, ready to reprice",
         size=16, color=DIM)

# Formula box
rounded_box(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.0),
            fill=PANEL_DARK, line=BLUE)
add_text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.0),
         "rank = (comp_median − your_price) × (1 + comp_count / 20) × confidence",
         size=18, bold=True, color=GREEN, font="Consolas",
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

bullets = [
    ("•  Ranks every listing by weighted upside × comp density × confidence", False, FG, 16),
    ("•  One-click deep link straight into Swipe Mode on the opportunity item", False, FG, 16),
    ("•  Filter by artist, min confidence, min upside %", False, FG, 16),
]
add_rich_lines(s, Inches(0.9), Inches(3.85), Inches(11.5), Inches(2.0),
               bullets, line_spacing=1.4)

# Big tag line
rounded_box(s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.95),
            fill=RGBColor(0x15, 0x24, 0x3A))
add_text(s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.95),
         "Stop clicking 200 listings. Start with the 10 that matter.",
         size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
wordmark(s)

# =========================================================================
# Slide 9 — Bulk Consensus Reprice + Drift Alerts (NEW)
# =========================================================================
s = new_slide()
title(s, "One click, many prices", accent=BLUE)

# Two columns
col_w = Inches(5.95)
col_h = Inches(4.5)
col_y = Inches(1.85)

# LEFT — Bulk Reprice
rounded_box(s, Inches(0.5), col_y, col_w, col_h, fill=PANEL, line=BLUE)
add_text(s, Inches(0.75), col_y + Inches(0.25), col_w, Inches(0.5),
         "BULK REPRICE MODAL", size=15, bold=True, color=BLUE)
bulk = [
    ("•  Filter: min upside %, min comp count, artist", False, FG, 14),
    ("•  Checkbox table → select-many → apply-many", False, FG, 14),
    ("•  NEW: \"Re-query LLMs\" checkbox forces fresh consensus", False, FG, 14),
    ("•  Source tag logged per change (manual, bulk_consensus, match_median)", False, FG, 14),
    ("•  Every change writes to the pricing track record", False, FG, 14),
]
add_rich_lines(s, Inches(0.75), col_y + Inches(0.85),
               col_w - Inches(0.4), col_h - Inches(1.0),
               bulk, line_spacing=1.45)

# RIGHT — Drift Alerts
rounded_box(s, Inches(6.85), col_y, col_w, col_h, fill=PANEL, line=GOLD)
add_text(s, Inches(7.1), col_y + Inches(0.25), col_w, Inches(0.5),
         "DRIFT ALERTS", size=15, bold=True, color=GOLD)
drift = [
    ("•  /api/drift-alerts surfaces comp-vs-listed gaps", False, FG, 14),
    ("•  Trigger: comp_median drifts ≥10% from your listed price", False, FG, 14),
    ("•  Catches missed windows around key-date events", False, FG, 14),
    ("•  Deep-links into Swipe Mode for one-click reprice", False, FG, 14),
    ("•  Morning heads-up on what moved overnight", False, FG, 14),
]
add_rich_lines(s, Inches(7.1), col_y + Inches(0.85),
               col_w - Inches(0.4), col_h - Inches(1.0),
               drift, line_spacing=1.45)

add_text(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.45),
         "Your pricing track record: every change, timestamped, sourced, reversible",
         size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
wordmark(s)

# =========================================================================
# Slide 10 — Comp Curation / Training Loop (NEW)
# =========================================================================
s = new_slide()
title(s, "DATARADAR gets smarter with every swipe", accent=GOLD)

bullets = [
    ("•  Individual comps surface in a swipe stack — reject noise, approve relevant", False, FG, 17),
    ("•  Rejections persist to comp_curation_rejections.json, keyed by sha1(name|price|date)", False, FG, 17),
    ("•  lookup_historical_prices auto-filters rejected comps → cleaner median / p75", False, FG, 17),
    ("•  Training data is the operator's expert eye, not labels at scale", False, FG, 17),
    ("•  Every rejection permanently tightens the model", False, FG, 17),
]
add_rich_lines(s, Inches(0.9), Inches(1.85), Inches(12), Inches(3.8),
               bullets, line_spacing=1.5)

# Visual: before / after comp distribution hint
row_y = Inches(5.55)
row_h = Inches(1.4)
# BEFORE
rounded_box(s, Inches(1.3), row_y, Inches(5.0), row_h,
            fill=PANEL, line=DIM)
add_text(s, Inches(1.3), row_y + Inches(0.1), Inches(5.0), Inches(0.4),
         "BEFORE CURATION", size=12, bold=True, color=DIM, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.3), row_y + Inches(0.5), Inches(5.0), Inches(0.9),
         "noisy comps  →  p75 wobble ±$40",
         size=17, bold=True, color=DIM, font="Consolas",
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# AFTER
rounded_box(s, Inches(7.1), row_y, Inches(5.0), row_h,
            fill=PANEL, line=GREEN)
add_text(s, Inches(7.1), row_y + Inches(0.1), Inches(5.0), Inches(0.4),
         "AFTER CURATION", size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, Inches(7.1), row_y + Inches(0.5), Inches(5.0), Inches(0.9),
         "clean comps  →  p75 tight ±$8",
         size=17, bold=True, color=GREEN, font="Consolas",
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Arrow between
arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                         Inches(6.35), row_y + Inches(0.55),
                         Inches(0.65), Inches(0.3))
arr.line.fill.background()
arr.fill.solid()
arr.fill.fore_color.rgb = GOLD
wordmark(s)

# =========================================================================
# Slide 11 — Architecture
# =========================================================================
s = new_slide()
title(s, "Under the hood", accent=GOLD)

# Center Flask box
center_w = Inches(3.3)
center_h = Inches(1.5)
cx = (SW - center_w) / 2
cy = Inches(3.3)
rounded_box(s, cx, cy, center_w, center_h, fill=PANEL, line=BLUE)
add_text(s, cx, cy, center_w, center_h,
         "Flask app\napp.py · 15,500 lines",
         size=15, bold=True, color=FG, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# Inputs (left side)
in_specs = [
    ("eBay Trading /\nBrowse API", Inches(0.5), Inches(1.65)),
    ("Google Sheets\n(pricing rules)", Inches(0.5), Inches(3.15)),
    ("54k comp DB\n(data/*.json)", Inches(0.5), Inches(4.65)),
]
for label, x, y in in_specs:
    rounded_box(s, x, y, Inches(2.5), Inches(1.1), fill=PANEL, line=GOLD)
    add_text(s, x, y, Inches(2.5), Inches(1.1),
             label, size=12, bold=True, color=FG, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               x + Inches(2.55), y + Inches(0.45),
                               cx - (x + Inches(2.55)) - Emu(30000),
                               Inches(0.22))
    arrow.line.fill.background()
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD

# Outputs (right side) — 4 LLM APIs in ThreadPoolExecutor
out_specs = [
    ("Claude API", Inches(10.3), Inches(1.4)),
    ("OpenAI API", Inches(10.3), Inches(2.65)),
    ("Gemini API", Inches(10.3), Inches(3.9)),
    ("Grok API", Inches(10.3), Inches(5.15)),
]
for label, x, y in out_specs:
    rounded_box(s, x, y, Inches(2.5), Inches(0.85), fill=PANEL, line=GREEN)
    add_text(s, x, y, Inches(2.5), Inches(0.85),
             label, size=12, bold=True, color=FG, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               cx + center_w + Emu(30000),
                               y + Inches(0.32),
                               x - (cx + center_w) - Emu(60000),
                               Inches(0.22))
    arrow.line.fill.background()
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GREEN

add_text(s, Inches(4.8), Inches(5.1), Inches(3.7), Inches(0.4),
         "parallel fan-out · ThreadPoolExecutor",
         size=11, color=DIM, align=PP_ALIGN.CENTER)

# Footer
add_text(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
         "Deployment: Railway + Cron (nightly reindex)  ·  "
         "Persistence: JSON + Google Sheets + Railway env vars",
         size=12, color=DIM, align=PP_ALIGN.CENTER)
wordmark(s)

# =========================================================================
# Slide 12 — Traction snapshot
# =========================================================================
s = new_slide()
title(s, "Where we are", accent=GREEN)

stats = [
    ("54,000+", "records in comp DB", GOLD),
    ("~200", "active eBay listings\nunder management", BLUE),
    ("3 of 4", "LLMs live\nClaude · GPT-4o · Gemini\nGrok pending env var", GREEN),
    ("<2s", "wall clock\n3-LLM fan-out", BLUE),
]
card_w = Inches(2.85)
card_h = Inches(3.2)
total = card_w * 4 + Inches(0.3) * 3
start = (SW - total) / 2
cy = Inches(2.0)
for i, (big_n, label, col) in enumerate(stats):
    x = start + (card_w + Inches(0.3)) * i
    rounded_box(s, x, cy, card_w, card_h, fill=PANEL, line=col)
    add_text(s, x, cy + Inches(0.35), card_w, Inches(1.4),
             big_n, size=50, bold=True, color=col, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), cy + Inches(1.85), card_w - Inches(0.4), Inches(1.3),
             label, size=13, color=FG, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.45),
         "Deployed live at web-production-15df7.up.railway.app",
         size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.4),
         "Railway auto-deploy  ·  Nightly comp re-index cron",
         size=13, color=DIM, align=PP_ALIGN.CENTER)
wordmark(s)

# =========================================================================
# Slide 13 — Strategic Options
# =========================================================================
s = new_slide()
title(s, "Strategic options", accent=GOLD)

opts = [
    ("OWN", GREEN,
     "Private moat driving Gauntlet Gallery margin — no one else has this pricing stack."),
    ("LICENSE", BLUE,
     "$99–299/mo SaaS for other high-volume art resellers  ·  100 subs × $99 ≈ $119k ARR."),
    ("SELL", GOLD,
     "Acquire-hire / IP sale to Heritage, 1stDibs, StockX, Rally.Rd, or eBay analytics."),
]
row_w = Inches(12.0)
row_h = Inches(1.3)
row_x = (SW - row_w) / 2
for i, (name, col, body) in enumerate(opts):
    y = Inches(1.85) + (row_h + Inches(0.25)) * i
    rounded_box(s, row_x, y, row_w, row_h, fill=PANEL, line=col)
    add_text(s, row_x + Inches(0.4), y, Inches(2.6), row_h,
             name, size=24, bold=True, color=col,
             anchor=MSO_ANCHOR.MIDDLE)
    sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             row_x + Inches(3.05), y + Inches(0.25),
                             Inches(0.03), row_h - Inches(0.5))
    sep.line.fill.background()
    sep.fill.solid()
    sep.fill.fore_color.rgb = col
    add_text(s, row_x + Inches(3.3), y, row_w - Inches(3.6), row_h,
             body, size=15, color=FG,
             anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.45),
         "Collectibles software TAM ≈ $400M  ·  growing 12% YoY",
         size=13, color=DIM, align=PP_ALIGN.CENTER)
wordmark(s)

# =========================================================================
# Slide 14 — Contact / End card
# =========================================================================
s = new_slide()
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                         Inches(5.5), Inches(1.5),
                         Inches(2.3), Inches(0.08))
bar.line.fill.background()
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE

add_text(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(0.9),
         "JJ Shay  ·  Gauntlet Gallery",
         size=40, bold=True, color=FG, align=PP_ALIGN.CENTER)

contact_lines = [
    ("jjshay@gmail.com", False, FG, 20),
    ("linkedin.com/in/jjshay", False, BLUE, 20),
    ("github.com/jjshay/dataradar-listings", False, GOLD, 20),
]
add_rich_lines(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(2.5),
               contact_lines, align=PP_ALIGN.CENTER, line_spacing=1.8)

add_text(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.6),
         "Thanks.",
         size=18, color=DIM, align=PP_ALIGN.CENTER)
wordmark(s)

# Save
out = "/tmp/dataradar-listings/DATARADAR_Strategic_Overview.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
